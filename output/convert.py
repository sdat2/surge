"""python output/convert.py"""
import os, sys
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

pdf_file = "output/presentation.pdf" # sys.argv[1] %
print()
print("Converting file: " + pdf_file)
print()

# Prep presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Create working folder
base_name = pdf_file.split(".pdf")[0]

# Convert PDF to list of images
print("Starting conversion...")
slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2)
print("...complete.")
print()

# Loop over slides
for i, slideimg in enumerate(slideimgs):
	if i % 3 == 0:
		print("Saving slide: " + str(i))

	imagefile = BytesIO()
	slideimg.save(imagefile, format='tiff')
	imagedata = imagefile.getvalue()
	imagefile.seek(0)
	width, height = slideimg.size

	# Set slide dimensions
	prs.slide_height = height * 9525
	prs.slide_width = width * 9525

	# Add slide
	slide = prs.slides.add_slide(blank_slide_layout)
	pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)
	if i == 4:
		video = os.path.join("../rotunno87", "gifs", "ani_model.mp4")
		image = os.path.join("images/phd/rot-day18.png")
		left = prs.slide_width*0.19
		top = prs.slide_height*0.11
		vwidth = prs.slide_width*0.62
		vheight = prs.slide_height*0.8
		shape = slide.shapes.add_shape(
			MSO_SHAPE.ROUNDED_RECTANGLE, left, top, vwidth, vheight
		)
		fill = shape.fill
		fill.solid()
		fill.fore_color.rgb = RGBColor(255, 255, 255)
		slide.shapes.add_movie(video, left, top, 
							   vwidth, vheight, poster_frame_image=image,
							   mime_type="video/mp4")
		
		#poster_frame_image=thumbnail_image_file)

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + '.pptx')
print("Conversion complete. :)")
print()

# slide.shapes.add_movie(video_file, x_pos, y_pos, width, height, poster_frame_image=thumbnail_image_file)
# slide_no=2
